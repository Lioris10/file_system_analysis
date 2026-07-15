"""Optional readers for PDF and Office documents."""

from __future__ import annotations

from pathlib import Path

from file_system_analysis.domain.exceptions import FileReadError


class PdfReader:
    """Read PDF text through pypdf when installed."""

    supported_extensions = {".pdf"}

    def read(self, path: Path) -> str:
        try:
            from pypdf import PdfReader as PypdfReader
        except ImportError as exc:
            raise FileReadError("pypdf is required to read PDF files") from exc

        reader = PypdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)


class DocxReader:
    """Read DOCX text through python-docx when installed."""

    supported_extensions = {".docx"}

    def read(self, path: Path) -> str:
        try:
            import docx
        except ImportError as exc:
            raise FileReadError("python-docx is required to read DOCX files") from exc

        document = docx.Document(str(path))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        for table in document.tables:
            for row in table.rows:
                paragraphs.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(paragraphs)


class PowerPointReader:
    """Read PPTX text through python-pptx when installed."""

    supported_extensions = {".pptx"}

    def read(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise FileReadError("python-pptx is required to read PPTX files") from exc

        presentation = Presentation(str(path))
        text_parts: list[str] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text_parts.append(f"Slide {slide_number}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
        return "\n".join(text_parts)


class SpreadsheetReader:
    """Read XLS/XLSX files through pandas when installed."""

    supported_extensions = {".xls", ".xlsx"}

    def read(self, path: Path) -> str:
        try:
            import pandas as pd
        except ImportError as exc:
            raise FileReadError("pandas with openpyxl/xlrd is required to read Excel files") from exc

        sheets = pd.read_excel(path, sheet_name=None, nrows=200)
        parts: list[str] = []
        for sheet_name, frame in sheets.items():
            parts.append(f"Sheet: {sheet_name}")
            parts.append(frame.to_csv(index=False))
        return "\n".join(parts)


class UnsupportedLegacyOfficeReader:
    """Report unsupported legacy binary Office formats in the MVP."""

    supported_extensions = {".doc", ".ppt"}

    def read(self, path: Path) -> str:
        raise FileReadError(f"Legacy Office format requires a conversion backend before reading: {path.suffix}")
